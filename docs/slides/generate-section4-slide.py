#!/usr/bin/env python3
"""Generate Section 4 PowerPoint slide: SPIFFE-verified VLAN management architecture."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Red Hat inspired palette
RH_RED = RGBColor(0xCC, 0x00, 0x00)
RH_DARK = RGBColor(0x21, 0x21, 0x21)
RH_BLUE = RGBColor(0x00, 0x66, 0xCC)
RH_GREEN = RGBColor(0x3E, 0x82, 0x27)
RH_ORANGE = RGBColor(0xEC, 0x7A, 0x08)
RH_GREY = RGBColor(0xF0, 0xF0, 0xF0)
RH_MID_GREY = RGBColor(0xA0, 0xA0, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_RED = RGBColor(0xFF, 0xE8, 0xE8)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF8, 0xE8)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
DARK_GREEN = RGBColor(0x2D, 0x6B, 0x1E)
DENY_RED = RGBColor(0xCC, 0x00, 0x00)
ALLOW_GREEN = RGBColor(0x3E, 0x82, 0x27)


def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None, text="", font_size=10, font_color=RH_DARK, bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top, color=RH_MID_GREY, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_text_box(slide, left, top, width, height, text, font_size=9, font_color=RH_DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def add_multiline_text_box(slide, left, top, width, height, lines, font_size=8, font_color=RH_DARK, alignment=PP_ALIGN.LEFT):
    """lines is a list of (text, bold, color) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, (text, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color if color else font_color
        run.font.bold = bold
    return txBox


def build_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RH_DARK
    title_bar.line.fill.background()
    tf = title_bar.text_frame
    tf.margin_left = Pt(20)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Section 4"
    run.font.size = Pt(14)
    run.font.color.rgb = RH_RED
    run.font.bold = True
    run2 = p.add_run()
    run2.text = "  SPIFFE-Verified Network VLAN Management"
    run2.font.size = Pt(22)
    run2.font.color.rgb = WHITE
    run2.font.bold = True

    # Subtitle
    add_text_box(slide, Inches(0.3), Inches(0.75), Inches(12), Inches(0.35),
                 "Defence in Depth: Two independent policy rings verify the human AND the workload before any network change proceeds",
                 font_size=11, font_color=RH_MID_GREY, bold=False, alignment=PP_ALIGN.LEFT)

    # ═══════════════════════════════════════════════════════
    # FLOW DIAGRAM — left to right across the slide
    # ═══════════════════════════════════════════════════════

    row_y = Inches(1.45)
    box_h = Inches(1.1)
    arrow_color = RGBColor(0x88, 0x88, 0x88)

    # Step 1: User + AAP
    step1_x = Inches(0.3)
    step1_w = Inches(2.0)
    add_rounded_box(slide, step1_x, row_y, step1_w, box_h,
                    LIGHT_BLUE, RH_BLUE,
                    "AAP Controller", font_size=11, bold=True, font_color=RH_BLUE)
    add_text_box(slide, step1_x, row_y + box_h + Pt(2), step1_w, Inches(0.35),
                 "netadmin launches\n\"Configure VLAN\" template",
                 font_size=8, font_color=RH_DARK, alignment=PP_ALIGN.CENTER)

    # Arrow 1→2
    add_arrow(slide, step1_x + step1_w, row_y + box_h / 2,
              step1_x + step1_w + Inches(0.35), row_y + box_h / 2, arrow_color)

    # Step 2: Outer Ring (AAP Gateway OPA)
    step2_x = step1_x + step1_w + Inches(0.4)
    step2_w = Inches(2.2)
    s2 = add_rounded_box(slide, step2_x, row_y, step2_w, box_h,
                         LIGHT_ORANGE, RH_ORANGE,
                         "", font_size=10, bold=True, font_color=RH_ORANGE)
    tf = s2.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(2)
    r1 = p1.add_run()
    r1.text = "OUTER RING"
    r1.font.size = Pt(9)
    r1.font.color.rgb = RH_ORANGE
    r1.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "AAP Gateway Policy"
    r2.font.size = Pt(11)
    r2.font.color.rgb = RH_DARK
    r2.font.bold = True
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(4)
    r3 = p3.add_run()
    r3.text = "aap.gateway OPA policy"
    r3.font.size = Pt(8)
    r3.font.color.rgb = RH_MID_GREY

    add_text_box(slide, step2_x, row_y + box_h + Pt(2), step2_w, Inches(0.35),
                 "Is user in Infrastructure\nAAP team?",
                 font_size=8, font_color=RH_DARK, alignment=PP_ALIGN.CENTER)

    # Arrow 2→3
    add_arrow(slide, step2_x + step2_w, row_y + box_h / 2,
              step2_x + step2_w + Inches(0.35), row_y + box_h / 2, arrow_color)

    # Step 3: SPIFFE verification
    step3_x = step2_x + step2_w + Inches(0.4)
    step3_w = Inches(2.2)
    s3 = add_rounded_box(slide, step3_x, row_y, step3_w, box_h,
                         LIGHT_GREEN, RH_GREEN,
                         "", font_size=10, bold=True)
    tf = s3.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(2)
    r1 = p1.add_run()
    r1.text = "SPIFFE / SPIRE"
    r1.font.size = Pt(11)
    r1.font.color.rgb = RH_GREEN
    r1.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Workload Identity"
    r2.font.size = Pt(10)
    r2.font.color.rgb = RH_DARK
    r2.font.bold = True
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(4)
    r3 = p3.add_run()
    r3.text = "X.509 SVID (short-lived)"
    r3.font.size = Pt(8)
    r3.font.color.rgb = RH_MID_GREY

    add_text_box(slide, step3_x, row_y + box_h + Pt(2), step3_w, Inches(0.35),
                 "Fetch SVID from SPIRE Agent\non AAP execution node",
                 font_size=8, font_color=RH_DARK, alignment=PP_ALIGN.CENTER)

    # Arrow 3→4
    add_arrow(slide, step3_x + step3_w, row_y + box_h / 2,
              step3_x + step3_w + Inches(0.35), row_y + box_h / 2, arrow_color)

    # Step 4: Inner Ring (OPA network policy)
    step4_x = step3_x + step3_w + Inches(0.4)
    step4_w = Inches(2.2)
    s4 = add_rounded_box(slide, step4_x, row_y, step4_w, box_h,
                         LIGHT_RED, RH_RED,
                         "", font_size=10, bold=True)
    tf = s4.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(2)
    r1 = p1.add_run()
    r1.text = "INNER RING"
    r1.font.size = Pt(9)
    r1.font.color.rgb = RH_RED
    r1.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "OPA Network Policy"
    r2.font.size = Pt(11)
    r2.font.color.rgb = RH_DARK
    r2.font.bold = True
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(4)
    r3 = p3.add_run()
    r3.text = "zta.network Rego policy"
    r3.font.size = Pt(8)
    r3.font.color.rgb = RH_MID_GREY

    add_text_box(slide, step4_x, row_y + box_h + Pt(2), step4_w, Inches(0.55),
                 "✓ Trusted SPIFFE ID?\n✓ User in network-admins?\n✓ VLAN 100-999?",
                 font_size=8, font_color=RH_DARK, alignment=PP_ALIGN.CENTER)

    # Arrow 4→5
    add_arrow(slide, step4_x + step4_w, row_y + box_h / 2,
              step4_x + step4_w + Inches(0.35), row_y + box_h / 2, arrow_color)

    # Step 5: Network + CMDB
    step5_x = step4_x + step4_w + Inches(0.4)
    step5_w = Inches(2.6)
    step5_h = Inches(0.45)

    add_rounded_box(slide, step5_x, row_y, step5_w, step5_h,
                    LIGHT_GREEN, DARK_GREEN,
                    "Arista cEOS Switches", font_size=10, bold=True, font_color=DARK_GREEN)
    add_rounded_box(slide, step5_x, row_y + step5_h + Inches(0.15), step5_w, step5_h,
                    LIGHT_BLUE, RH_BLUE,
                    "NetBox CMDB Update", font_size=10, bold=True, font_color=RH_BLUE)

    add_text_box(slide, step5_x, row_y + box_h + Pt(2), step5_w, Inches(0.35),
                 "VLAN created + audit trail\nwith user + workload identity",
                 font_size=8, font_color=RH_DARK, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════
    # DENY paths — below the main flow
    # ═══════════════════════════════════════════════════════

    deny_y = Inches(3.3)

    # Outer ring deny
    deny1 = add_rounded_box(slide, step2_x + Inches(0.2), deny_y, Inches(1.8), Inches(0.45),
                            RGBColor(0xFF, 0xDD, 0xDD), DENY_RED,
                            "✗ DENIED", font_size=10, bold=True, font_color=DENY_RED)
    add_text_box(slide, step2_x - Inches(0.1), deny_y + Inches(0.48), Inches(2.4), Inches(0.3),
                 "neteng blocked — wrong team",
                 font_size=8, font_color=DENY_RED, alignment=PP_ALIGN.CENTER)
    add_arrow(slide, step2_x + step2_w / 2, row_y + box_h,
              step2_x + step2_w / 2, deny_y, DENY_RED, Pt(1.5))

    # Inner ring deny
    deny2 = add_rounded_box(slide, step4_x + Inches(0.2), deny_y, Inches(1.8), Inches(0.45),
                            RGBColor(0xFF, 0xDD, 0xDD), DENY_RED,
                            "✗ DENIED", font_size=10, bold=True, font_color=DENY_RED)
    add_text_box(slide, step4_x - Inches(0.1), deny_y + Inches(0.48), Inches(2.4), Inches(0.3),
                 "Rogue SPIFFE ID / bad VLAN / wrong group",
                 font_size=8, font_color=DENY_RED, alignment=PP_ALIGN.CENTER)
    add_arrow(slide, step4_x + step4_w / 2, row_y + box_h,
              step4_x + step4_w / 2, deny_y, DENY_RED, Pt(1.5))

    # ═══════════════════════════════════════════════════════
    # BOTTOM SECTION: Infrastructure components
    # ═══════════════════════════════════════════════════════

    infra_y = Inches(4.4)
    add_text_box(slide, Inches(0.3), infra_y - Inches(0.3), Inches(4), Inches(0.25),
                 "Lab Infrastructure", font_size=11, font_color=RH_DARK, bold=True)

    comp_h = Inches(0.7)
    comp_w = Inches(2.3)
    gap = Inches(0.3)

    components = [
        ("SPIRE Server\ncentral.zta.lab", "Trust domain CA\nSigns all SVIDs", LIGHT_GREEN, RH_GREEN),
        ("SPIRE Agent\ncontrol.zta.lab", "Serves SVIDs to AAP\nvia Unix socket", LIGHT_GREEN, RH_GREEN),
        ("OPA\ncentral.zta.lab:8181", "aap.gateway (outer)\nzta.network (inner)", LIGHT_ORANGE, RH_ORANGE),
        ("IdM (FreeIPA)\ncentral.zta.lab", "User identity + groups\nnetwork-admins", LIGHT_BLUE, RH_BLUE),
        ("NetBox\nnetbox.zta.lab:8880", "CMDB audit trail\nVLAN + identity records", LIGHT_BLUE, RH_BLUE),
    ]

    for i, (title, desc, bg_color, border) in enumerate(components):
        x = Inches(0.3) + i * (comp_w + gap)
        box = add_rounded_box(slide, x, infra_y, comp_w, comp_h,
                              bg_color, border, "", font_size=9)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(2)
        r1 = p1.add_run()
        title_parts = title.split("\n")
        r1.text = title_parts[0]
        r1.font.size = Pt(9)
        r1.font.color.rgb = border
        r1.font.bold = True
        if len(title_parts) > 1:
            p1b = tf.add_paragraph()
            p1b.alignment = PP_ALIGN.CENTER
            p1b.space_before = Pt(0)
            r1b = p1b.add_run()
            r1b.text = title_parts[1]
            r1b.font.size = Pt(7)
            r1b.font.color.rgb = RH_MID_GREY
        for line in desc.split("\n"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(1)
            r2 = p2.add_run()
            r2.text = line
            r2.font.size = Pt(7)
            r2.font.color.rgb = RH_DARK

    # ═══════════════════════════════════════════════════════
    # KEY TAKEAWAYS — bottom right
    # ═══════════════════════════════════════════════════════

    takeaway_x = Inches(0.3)
    takeaway_y = Inches(5.4)
    takeaway_w = Inches(12.7)
    takeaway_h = Inches(0.65)

    tb = add_rounded_box(slide, takeaway_x, takeaway_y, takeaway_w, takeaway_h,
                         RH_GREY, None, "", font_size=8)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.space_before = Pt(2)
    r0 = p0.add_run()
    r0.text = "Key Takeaways"
    r0.font.size = Pt(9)
    r0.font.color.rgb = RH_DARK
    r0.font.bold = True

    takeaways = [
        "Outer ring (AAP gateway) checks WHO can press the button  •  Inner ring (in-playbook OPA) checks WHAT the workload is and WHAT it requests",
        "SPIFFE SVIDs are short-lived and auto-rotated — stolen certificates expire in minutes, not months  •  Authentication (SPIRE) ≠ Authorisation (OPA)",
    ]
    for t in takeaways:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        r = p.add_run()
        r.text = "▸ " + t
        r.font.size = Pt(8)
        r.font.color.rgb = RH_DARK

    # Footer
    add_text_box(slide, Inches(0.3), Inches(7.1), Inches(6), Inches(0.25),
                 "Zero Trust Architecture Workshop  •  Ansible Automation Platform 2.6",
                 font_size=7, font_color=RH_MID_GREY)
    add_text_box(slide, Inches(9), Inches(7.1), Inches(4), Inches(0.25),
                 "Section 4: SPIFFE + OPA + Arista + NetBox",
                 font_size=7, font_color=RH_MID_GREY, alignment=PP_ALIGN.RIGHT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide(prs)

    output_path = "docs/slides/section4-architecture.pptx"
    prs.save(output_path)
    print(f"Saved to {output_path}")


if __name__ == "__main__":
    main()
